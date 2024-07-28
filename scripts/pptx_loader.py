from pptx import Presentation

def load_pptx(filename):
    # 打开PPT文件
    prs = Presentation(filename)
    retval_content = ""
    # 遍历所有的幻灯片
    for slide in prs.slides:
        # 提取幻灯片的标题
        title = ""
        if slide.shapes.title != None:
            title = slide.shapes.title.text.strip()
        # 提取幻灯片的正文文本
        content = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                content += shape.text.strip() + "\n"

        retval_content += title + "\n" + content
        # print("标题:", title)
        # print("内容:", content)
    contents_list = retval_content.strip().split("\n")
    retval = ""
    for item in contents_list:
        if item == "":
            continue
        retval += item + "\n"
    # print(retval)
    return retval


load_pptx("./2022年度述职PPT-资金清算处.pptx")